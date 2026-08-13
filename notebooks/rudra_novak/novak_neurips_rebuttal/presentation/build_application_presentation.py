#!/usr/bin/env python3
"""Build the NeuralFieldManifold rebuttal application presentation."""

from __future__ import annotations

import csv
import json
import math
import textwrap
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable, Literal, Sequence

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5
PREVIEW_W = 1600
PREVIEW_H = 900

WHITE = "FFFFFF"
INK = "191919"
CHARCOAL = "4B4B4B"
MUTED = "777777"
RED = "8F160D"
MID_RED = "C94C28"
LIGHT_RED = "F5E4E0"
PALE = "F5F3F2"
GRID = "DDD8D6"
GREEN = "287D5A"
BLUE = "315C85"
BLACK = "000000"

FONT_NAME = "Aptos"
PREVIEW_FONT = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
PREVIEW_FONT_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"


@dataclass(frozen=True)
class TextElement:
    x: float
    y: float
    w: float
    h: float
    text: str
    size: float = 18
    color: str = INK
    bold: bool = False
    align: Literal["left", "center", "right"] = "left"
    valign: Literal["top", "middle", "bottom"] = "top"
    fill: str | None = None
    line: str | None = None
    margin: float = 0.08


@dataclass(frozen=True)
class ImageElement:
    x: float
    y: float
    w: float
    h: float
    path: Path
    mode: Literal["fit", "crop"] = "fit"
    line: str | None = None


@dataclass(frozen=True)
class RectElement:
    x: float
    y: float
    w: float
    h: float
    fill: str = PALE
    line: str | None = None
    radius: float = 0.0


@dataclass(frozen=True)
class LineElement:
    x: float
    y: float
    w: float
    h: float
    color: str = RED
    width: float = 2.0


Element = TextElement | ImageElement | RectElement | LineElement


@dataclass(frozen=True)
class SlideSpec:
    title: str
    section: str
    kind: Literal["core", "backup"]
    elements: tuple[Element, ...]
    script: str
    takeaway: str
    transition: str
    challenge_response: str
    speaker_cue: str = ""

    def __post_init__(self) -> None:
        for element in self.elements:
            if (
                element.x < 0
                or element.y < 0
                or element.w < 0
                or element.h < 0
                or element.x + element.w > SLIDE_W + 1e-6
                or element.y + element.h > SLIDE_H + 1e-6
            ):
                raise ValueError(f"Element outside slide bounds: {element}")

    @property
    def text_elements(self) -> tuple[TextElement, ...]:
        return tuple(e for e in self.elements if isinstance(e, TextElement))


def _row(df: pd.DataFrame, **filters: Any) -> pd.Series:
    view = df
    for key, value in filters.items():
        view = view[view[key] == value]
    if len(view) != 1:
        raise ValueError(f"Expected one row for {filters}; found {len(view)}")
    return view.iloc[0]


def _metric(row: pd.Series) -> dict[str, float]:
    return {
        "mean_f1": float(row["mean_f1"]),
        "std_f1": float(row["std_f1"]),
    }


def load_results(root: Path) -> dict[str, Any]:
    tables = root / "tables"
    motor = root / "motor_lfp_reaching"
    motor_tables = motor / "tables"

    eeg_summary = pd.read_csv(tables / "eeg_reviewer_feature_summary.csv")
    eeg_sig = pd.read_csv(tables / "eeg_reviewer_relevantband_significance.csv")
    eeg_ablation = pd.read_csv(tables / "eeg_torus_order_ablation_summary.csv")
    eeg_ablation_sig = pd.read_csv(tables / "eeg_torus_order_ablation_significance.csv")
    session_counts = pd.read_csv(tables / "session_class_counts.csv")

    motor_summary = pd.read_csv(
        motor_tables / "nonlinear_refit_direction_movement_summary_pertrace_tau_dim.csv"
    )
    motor_sig = pd.read_csv(
        motor_tables / "torus_avgpsd_relevantband_significance_pertrace_tau_dim.csv"
    )
    additive = pd.read_csv(
        motor_tables / "additive_torus_feature_summary_pertrace_tau_dim.csv"
    )
    additive_sig = pd.read_csv(
        motor_tables / "additive_torus_significance_pertrace_tau_dim.csv"
    )
    embedding = pd.read_csv(motor_tables / "lfp_embedding_params_pertrace_tau_dim.csv")
    manifest = pd.read_csv(motor_tables / "lfp_manifest.csv")
    by_monkey = pd.read_csv(
        motor_tables / "torus_avgpsd_relevantband_f1_bar_summary_by_monkey_pertrace_tau_dim.csv"
    )

    eeg = {
        "relevant": _metric(_row(eeg_summary, feature_set="delta")),
        "average_psd": _metric(_row(eeg_summary, feature_set="average_psd")),
        "all_band": _metric(_row(eeg_summary, feature_set="all_band_power")),
        "torus": _metric(_row(eeg_summary, feature_set="all_torus_15")),
    }
    eeg["torus_vs_relevant_p_holm"] = float(
        _row(eeg_sig, comparison="torus_features vs relevant_band")["p_holm"]
    )
    eeg["torus_vs_psd_p_holm"] = float(
        _row(eeg_sig, comparison="torus_features vs average_psd")["p_holm"]
    )
    eeg["ablation_full"] = _metric(_row(eeg_ablation, feature_set="torus_features_15d"))
    eeg["ablation_no_r"] = _metric(_row(eeg_ablation, feature_set="no_r_geometry_11d"))
    eeg["ablation_p"] = float(eeg_ablation_sig.iloc[0]["p_value"])

    full_motor = motor_summary[motor_summary["analysis_set"] == "full_6_direction_lfps"]
    motor_results = {
        "relevant": _metric(_row(full_motor, feature_set="beta")),
        "average_psd": _metric(_row(full_motor, feature_set="average_psd")),
        "all_band": _metric(_row(full_motor, feature_set="all_band_power")),
        "torus": _metric(_row(full_motor, feature_set="torus_nonlinear_15")),
    }
    motor_results["torus_vs_relevant_p_holm"] = float(
        _row(motor_sig, monkey="M+T", comparison="torus_features vs relevant_band")["p_holm"]
    )
    motor_results["torus_vs_psd_p_holm"] = float(
        _row(motor_sig, monkey="M+T", comparison="torus_features vs average_psd")["p_holm"]
    )

    pooled_add = additive[additive["analysis_level"] == "pooled"]
    for feature in [
        "torus_plus_relevant_band",
        "torus_plus_average_psd",
        "torus_plus_all_band_power",
    ]:
        motor_results[feature] = _metric(_row(pooled_add, feature_set=feature))

    primary_add = additive_sig[
        (additive_sig["analysis_level"] == "pooled")
        & (additive_sig["comparison_family"] == "primary")
    ]
    motor_results["additive_significance"] = {
        row["comparison"]: {
            "p_holm": float(row["p_holm"]),
            "difference": float(row["mean_difference_left_minus_right"]),
        }
        for _, row in primary_add.iterrows()
    }

    embedding_ok = embedding[embedding["status"] == "ok"].copy()
    dims = embedding_ok["torus_embedding_dim"].astype(int).value_counts().sort_index()

    results = {
        "root": root,
        "motor_root": motor,
        "eeg": eeg,
        "motor": motor_results,
        "embedding": {
            "median_tau_ms": float(embedding_ok["torus_tau_ms"].median()),
            "tau_q1_ms": float(embedding_ok["torus_tau_ms"].quantile(0.25)),
            "tau_q3_ms": float(embedding_ok["torus_tau_ms"].quantile(0.75)),
            "dim_counts": {int(k): int(v) for k, v in dims.items()},
        },
        "by_monkey": by_monkey,
        "counts": {
            "eeg_total_sessions": int(len(session_counts)),
            "eeg_valid_sessions": int((session_counts["status"] == "ok").sum()),
            "motor_unique_lfps": int(len(manifest)),
            "motor_full_lfps": int(manifest["has_all_6_directions"].astype(bool).sum()),
        },
    }
    return results


def validate_claims(results: dict[str, Any]) -> None:
    expected = {
        ("eeg", "torus", "mean_f1"): 0.6697563098214888,
        ("eeg", "all_band", "mean_f1"): 0.695717725985247,
        ("motor", "torus", "mean_f1"): 0.18062425168818425,
        ("motor", "all_band", "mean_f1"): 0.21157628209290785,
    }
    for path, target in expected.items():
        value: Any = results
        for key in path:
            value = value[key]
        if not math.isclose(float(value), target, rel_tol=0, abs_tol=1e-10):
            raise ValueError(f"Numerical claim changed for {path}: {value} != {target}")
    if results["counts"] != {
        "eeg_total_sessions": 24,
        "eeg_valid_sessions": 21,
        "motor_unique_lfps": 341,
        "motor_full_lfps": 237,
    }:
        raise ValueError(f"Unexpected dataset counts: {results['counts']}")


def load_eeg_torus_confusion_summary(root: Path) -> dict[str, Any]:
    matrices: list[np.ndarray] = []
    session_ids: list[str] = []
    for path in sorted((root / "cache").glob("session_*.npz")):
        with np.load(path, allow_pickle=True) as data:
            feature_sets = [str(value) for value in data["feature_sets"]]
            if "all_torus_15" not in feature_sets:
                continue
            matrix = np.asarray(
                data["confusion_mats"][feature_sets.index("all_torus_15")], dtype=float
            )
        if matrix.shape == (3, 3) and np.isfinite(matrix).all():
            matrices.append(matrix)
            session_ids.append(path.stem)
    if not matrices:
        raise ValueError("No valid EEG torus confusion matrices were found")
    stacked = np.stack(matrices)
    return {
        "n_sessions": len(matrices),
        "session_ids": session_ids,
        "mean": stacked.mean(axis=0),
        "std": stacked.std(axis=0, ddof=1),
    }


def render_eeg_class_accuracy_figure(root: Path, output: Path) -> None:
    summary = load_eeg_torus_confusion_summary(root)
    original = np.array(
        [
            [0.65, 0.06, 0.29],
            [0.02, 0.94, 0.05],
            [0.32, 0.05, 0.63],
        ],
        dtype=float,
    )
    labels = ["Wake", "NREM", "REM"]
    cmap = LinearSegmentedColormap.from_list(
        "paper_red", ["#f5f3f2", "#d45f57", "#8f160d"]
    )
    fig, axes = plt.subplots(1, 2, figsize=(12.2, 5.0), constrained_layout=True)
    panels = [
        (axes[0], original, None, "Original example"),
        (
            axes[1],
            np.asarray(summary["mean"]),
            np.asarray(summary["std"]),
            f"Across {summary['n_sessions']} recording hours",
        ),
    ]
    image = None
    for ax, matrix, variability, title in panels:
        image = ax.imshow(matrix, vmin=0, vmax=1, cmap=cmap)
        ax.set_xticks(range(3), labels)
        ax.set_yticks(range(3), labels)
        ax.set_xlabel("Predicted stage", fontsize=12)
        ax.set_ylabel("True stage", fontsize=12)
        ax.set_title(title, fontsize=15, fontweight="bold", pad=10)
        for row in range(3):
            for col in range(3):
                value = matrix[row, col]
                annotation = f"{value:.2f}"
                if variability is not None:
                    annotation = f"{value:.2f}\n+/- {variability[row, col]:.2f}"
                color = "white" if value >= 0.52 else "#191919"
                ax.text(
                    col,
                    row,
                    annotation,
                    ha="center",
                    va="center",
                    fontsize=13 if variability is None else 11,
                    fontweight="bold" if row == col else "normal",
                    color=color,
                )
        for spine in ax.spines.values():
            spine.set_linewidth(0.8)
    assert image is not None
    colorbar = fig.colorbar(image, ax=axes, shrink=0.82, pad=0.035)
    colorbar.set_label("Row-normalized decoding accuracy", fontsize=11)
    output.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def fmt_f1(metric: dict[str, float]) -> str:
    return f"{metric['mean_f1']:.3f} +/- {metric['std_f1']:.3f} SD"


def fmt_p(value: float) -> str:
    if value < 0.001:
        exponent = int(math.floor(math.log10(value)))
        mantissa = value / (10**exponent)
        return f"p = {mantissa:.2f} x 10^{exponent}"
    return f"p = {value:.3f}"


def t(
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align: Literal["left", "center", "right"] = "left",
    valign: Literal["top", "middle", "bottom"] = "top",
    fill: str | None = None,
    line: str | None = None,
    margin: float = 0.08,
) -> TextElement:
    return TextElement(x, y, w, h, text, size, color, bold, align, valign, fill, line, margin)


def img(x: float, y: float, w: float, h: float, path: Path, mode: str = "fit") -> ImageElement:
    return ImageElement(x, y, w, h, path, mode)  # type: ignore[arg-type]


def rect(x: float, y: float, w: float, h: float, fill: str, line: str | None = None) -> RectElement:
    return RectElement(x, y, w, h, fill, line)


def line(x: float, y: float, w: float, h: float, color: str = RED, width: float = 2.0) -> LineElement:
    return LineElement(x, y, w, h, color, width)


def metric_card(x: float, y: float, label: str, value: str, accent: str = RED, w: float = 3.55) -> list[Element]:
    return [
        rect(x, y, w, 1.18, WHITE, GRID),
        rect(x, y, 0.09, 1.18, accent),
        t(x + 0.23, y + 0.15, w - 0.35, 0.28, label.upper(), 11, MUTED, True),
        t(x + 0.23, y + 0.47, w - 0.35, 0.52, value, 22, INK, True),
    ]


def bullet_block(x: float, y: float, w: float, items: Sequence[str], size: float = 18, color: str = INK) -> list[Element]:
    elements: list[Element] = []
    row_h = 0.62 if size >= 18 else 0.52
    for index, item in enumerate(items):
        cy = y + index * row_h
        elements.append(t(x, cy + 0.05, 0.28, 0.30, "-", size, RED, True, "center"))
        elements.append(t(x + 0.34, cy, w - 0.34, row_h, item, size, color))
    return elements


def standard_slide(
    title: str,
    section: str,
    kind: Literal["core", "backup"],
    body: Sequence[Element],
    script: str,
    takeaway: str,
    transition: str,
    challenge: str,
    cue: str = "",
) -> SlideSpec:
    chrome: list[Element] = [
        t(0.55, 0.18, 4.2, 0.24, "BACKUP" if kind == "backup" else section.upper(), 10, RED, True),
        t(0.55, 0.46, 12.15, 0.52, title, 27, INK, True),
        line(0.55, 1.07, 12.22, 0.0, RED, 2.0),
    ]
    return SlideSpec(title, section, kind, tuple(chrome + list(body)), script, takeaway, transition, challenge, cue)


def title_slide(title: str, subtitle: str, script: str) -> SlideSpec:
    body: list[Element] = [
        rect(0, 0, SLIDE_W, SLIDE_H, WHITE),
        rect(0, 0, 0.22, SLIDE_H, RED),
        t(0.85, 0.70, 3.0, 0.28, "NEURALFIELDMANIFOLD", 12, RED, True),
        t(0.85, 1.55, 11.4, 1.35, title, 38, INK, True),
        t(0.87, 3.05, 10.7, 0.85, subtitle, 22, CHARCOAL),
        line(0.87, 4.28, 2.3, 0, RED, 4),
        t(0.87, 4.60, 10.8, 0.7, "Mouse EEG sleep staging  |  Macaque reach-direction decoding", 17, MUTED),
        t(0.87, 6.68, 9.8, 0.35, "NeurIPS rebuttal experiments - application section", 12, MUTED),
    ]
    return SlideSpec(
        title,
        "Applications",
        "core",
        tuple(body),
        script,
        "The application section asks whether the recovered geometry carries behaviorally useful information.",
        "I will start with the concrete questions raised during review.",
        "If asked why decoding matters: it is a proof-of-concept behavioral readout, not the primary claim of the paper.",
        "Pause after the title, then explicitly connect to Kasra's theory section.",
    )


def build_full_slide_specs(results: dict[str, Any]) -> list[SlideSpec]:
    root = Path(results["root"])
    motor = Path(results["motor_root"])
    eeg_plots = root / "plots" / "summary"
    motor_plots = motor / "plots" / "summary"

    eeg = results["eeg"]
    mot = results["motor"]
    counts = results["counts"]
    emb = results["embedding"]

    slides: list[SlideSpec] = []

    slides.append(
        title_slide(
            "From torus geometry to behavior",
            "What the rebuttal experiments tell us - and what they do not",
            "Kasra has shown the theoretical connection between sustained oscillatory modes and toroidal lag geometry. My part asks the next empirical question: once we recover that geometry, does it carry information about what the brain and animal are doing? The rebuttal gave us two tests: sleep staging across a full day of mouse EEG and six-way reach-direction decoding from macaque motor-cortex LFP.",
        )
    )

    reviewer_cards: list[Element] = []
    questions = [
        ("BASELINES", "Does the decoding gap survive explicit spectral baselines, cross-validation, and error bars?"),
        ("GENERALIZATION", "Is the result stable across recording hours and separately in each animal?"),
        ("NEW TASK", "Can the method decode repeated reaching movements from motor-cortex LFP?"),
    ]
    for i, (label, text_value) in enumerate(questions):
        x = 0.72 + i * 4.15
        reviewer_cards.extend(
            [
                rect(x, 1.55, 3.72, 3.65, WHITE, GRID),
                t(x + 0.28, 1.86, 3.15, 0.32, f"0{i + 1}", 22, RED, True),
                t(x + 0.28, 2.38, 3.10, 0.30, label, 12, MUTED, True),
                t(x + 0.28, 2.86, 3.08, 1.72, text_value, 20, INK, True),
            ]
        )
    reviewer_cards.append(
        t(0.78, 5.62, 11.9, 0.72, "The rebuttal became a stress test of behavioral utility, reproducibility, and claim precision.", 21, RED, True, "center")
    )
    slides.append(
        standard_slide(
            "Reviewers asked for evidence beyond a single decoding panel",
            "Motivation",
            "core",
            reviewer_cards,
            "The reviewers did not simply ask for a larger number. They asked whether the result survives a fully specified baseline, whether it is stable across independent recording units, and whether the same idea works in a repeated motor task. Those requests shaped every analysis I will show.",
            "The new experiments test rigor and scope, not just performance.",
            "First, here is the common evaluation framework used across both datasets.",
            "If asked which reviewer prompted what, use the backup evidence matrix; the live story is organized by scientific question.",
        )
    )

    pipeline: list[Element] = []
    boxes = [
        ("WINDOW", "Task-aligned or 2 s signal segment"),
        ("FEATURES", "Single band | Average PSD | All bands | Torus"),
        ("DECODER", "Balanced, stratified 5-fold LDA"),
        ("REPORT", "Macro F1; SD across sessions or LFPs"),
    ]
    for i, (label, text_value) in enumerate(boxes):
        x = 0.55 + i * 3.16
        pipeline.extend(
            [
                rect(x, 1.65, 2.67, 1.42, WHITE, GRID),
                t(x + 0.18, 1.88, 2.3, 0.25, label, 11, RED, True),
                t(x + 0.18, 2.25, 2.28, 0.58, text_value, 16, INK, True),
            ]
        )
        if i < 3:
            pipeline.extend([line(x + 2.68, 2.35, 0.45, 0, RED, 2.3), t(x + 2.86, 2.15, 0.25, 0.28, ">", 18, RED, True)])
    pipeline.extend(
        metric_card(0.75, 4.18, "Statistical unit", "Matched session or LFP", BLUE, 3.72)
        + metric_card(4.80, 4.18, "Test", "Paired Wilcoxon", RED, 3.72)
        + metric_card(8.85, 4.18, "Multiplicity", "Holm correction", MID_RED, 3.72)
    )
    pipeline.append(t(0.76, 5.78, 11.75, 0.55, "Chance: 0.33 for three sleep stages | 0.17 for six reach directions", 17, MUTED, True, "center"))
    slides.append(
        standard_slide(
            "Both datasets use the same conservative decoding logic",
            "Methods",
            "core",
            pipeline,
            "For every comparison, we held the decoder fixed and only changed the representation. We balanced the labels, used stratified five-fold LDA, and report macro F1. Error bars are standard deviation across independent sessions for EEG or independent LFP traces for the macaque. Statistical comparisons are paired Wilcoxon tests with Holm correction.",
            "Differences in performance reflect representations, not different classifiers.",
            "With that framework fixed, the first test is the full-day mouse EEG recording.",
            "Macro F1 is appropriate because it weights Wake, NREM, and REM equally after balancing and remains interpretable across six reach directions.",
        )
    )

    eeg_overview = [
        rect(0.65, 1.45, 4.2, 4.7, PALE),
        t(0.95, 1.86, 3.55, 0.60, "24", 45, RED, True, "center"),
        t(0.95, 2.55, 3.55, 0.55, "one-hour sessions", 20, INK, True, "center"),
        t(0.95, 3.42, 3.55, 0.60, "21", 45, RED, True, "center"),
        t(0.95, 4.10, 3.55, 0.82, "sessions with enough Wake, NREM, and REM", 18, INK, True, "center"),
        t(0.95, 5.35, 3.55, 0.45, "2-second windows", 17, MUTED, True, "center"),
        t(5.35, 1.65, 6.95, 0.55, "Question", 13, RED, True),
        t(5.35, 2.16, 6.75, 1.00, "Can compact geometry decode sleep stage consistently across recording hours?", 25, INK, True),
    ] + bullet_block(
        5.38,
        3.55,
        6.7,
        [
            "Three classes: Wake, NREM, REM",
            "Balanced sampling within each hour",
            "Same feature definitions and classifier throughout",
            "Session is the unit of replication",
        ],
        18,
    )
    slides.append(
        standard_slide(
            "Mouse EEG provides a 24-hour repeated-session test",
            "Mouse EEG",
            "core",
            eeg_overview,
            "We converted the full 24-hour cortical EEG recording into 24 one-hour sessions. Twenty-one hours contained at least five windows from each sleep class and entered balanced decoding. This gives us repeated temporal replicates rather than one pooled result.",
            "The EEG analysis tests repeatability across hours, not just average performance.",
            "The hour-by-hour heatmap shows how stable the representations are.",
            "The three excluded hours lacked sufficient representation of all three classes; they were marked insufficient rather than forced into cross-validation.",
        )
    )

    slides.append(
        standard_slide(
            "Torus decoding remains stable across recording hours",
            "Mouse EEG",
            "core",
            [
                img(0.72, 1.30, 8.62, 5.72, eeg_plots / "session_f1_heatmap.png", "fit"),
                t(9.62, 1.65, 2.75, 0.35, "READOUT", 11, RED, True),
                t(9.62, 2.10, 2.75, 1.12, "Torus F1 spans approximately 0.57-0.75", 23, INK, True),
                t(9.62, 3.52, 2.75, 1.42, "Torus exceeds every single band in 17 of 21 held-out hours and ties the strongest band in two more.", 17, CHARCOAL),
                t(9.62, 5.40, 2.70, 0.72, "Variation is real, but the effect is not carried by one exceptional hour.", 16, RED, True),
            ],
            "Each row is one valid recording hour, and each column is a feature set. The important pattern is not a perfectly flat torus column; it is that torus decoding remains high across hours and is usually stronger than any single frequency band.",
            "The EEG result is distributed across time rather than driven by one session.",
            "Collapsing across hours lets us compare the compact baselines directly.",
            "These are within-hour cross-validated F1 values summarized across hours; they should not be described as a train-on-one-hour, test-on-another transfer experiment.",
        )
    )

    eeg_comp: list[Element] = []
    eeg_comp.extend(metric_card(0.65, 1.62, "Relevant band: delta", fmt_f1(eeg["relevant"]), MID_RED, 3.65))
    eeg_comp.extend(metric_card(4.84, 1.62, "Average PSD", fmt_f1(eeg["average_psd"]), "888888", 3.65))
    eeg_comp.extend(metric_card(9.02, 1.62, "Torus features", fmt_f1(eeg["torus"]), RED, 3.65))
    eeg_comp.extend(
        [
            line(2.35, 3.35, 8.05, 0, RED, 1.8),
            t(3.10, 3.47, 3.00, 0.42, fmt_p(eeg["torus_vs_relevant_p_holm"]), 16, RED, True, "center"),
            t(7.25, 3.47, 3.00, 0.42, fmt_p(eeg["torus_vs_psd_p_holm"]), 16, RED, True, "center"),
            t(0.95, 4.40, 11.45, 0.58, "Geometry improves F1 by 0.200 over delta and 0.185 over Average PSD.", 25, INK, True, "center"),
            t(1.15, 5.33, 11.05, 0.80, "This is the defensible reviewer-facing claim: torus geometry carries information beyond compact spectral summaries.", 19, CHARCOAL, False, "center"),
        ]
    )
    slides.append(
        standard_slide(
            "EEG geometry outperforms compact spectral summaries",
            "Mouse EEG",
            "core",
            eeg_comp,
            f"Across 21 sessions, torus features reach {fmt_f1(eeg['torus'])}. Delta, our physiologically relevant sleep band, reaches {fmt_f1(eeg['relevant'])}, and a single scalar Average PSD reaches {fmt_f1(eeg['average_psd'])}. Both paired differences survive Holm correction at {fmt_p(eeg['torus_vs_relevant_p_holm'])}.",
            "Torus geometry is substantially more informative than delta power or mean broadband PSD alone.",
            "The next slide is the qualification that keeps this claim honest.",
            "Relevant band is fixed a priori as delta for sleep; it is not selected after seeing decoder performance.",
        )
    )

    slides.append(
        standard_slide(
            "The joint all-band vector is stronger than torus alone",
            "Mouse EEG",
            "core",
            [
                img(0.60, 1.30, 8.20, 5.62, eeg_plots / "eeg_sleep_feature_f1_barplot.png", "fit"),
                rect(9.10, 1.68, 3.48, 3.02, LIGHT_RED, RED),
                t(9.42, 1.98, 2.85, 0.35, "IMPORTANT NUANCE", 12, RED, True),
                t(9.42, 2.55, 2.88, 0.67, fmt_f1(eeg["all_band"]), 24, INK, True),
                t(9.42, 3.38, 2.86, 0.94, "All-band power combines six band powers in one feature vector.", 17, CHARCOAL),
                t(9.18, 5.20, 3.28, 0.95, "Claim: geometry complements spectral summaries. Not: geometry universally beats spectral features.", 18, RED, True, "center"),
            ],
            f"When we give the spectral decoder all six band-power values jointly, it reaches {fmt_f1(eeg['all_band'])}, slightly above the torus result. That does not erase the geometry result; it sharpens it. The torus representation beats each single band and Average PSD, but a richer multiband spectral vector remains a strong baseline.",
            "The evidence supports complementary information, not universal spectral dominance.",
            "We also tested how much of the torus feature set is actually necessary.",
            "All-band power is a multifeature vector, not the average of the band powers. Its dimensionality and information content are larger than the one-scalar baselines.",
        )
    )

    slides.append(
        standard_slide(
            "Removing radius-related terms causes a modest, nonsignificant drop",
            "Mouse EEG",
            "core",
            [
                img(0.75, 1.25, 7.15, 5.85, eeg_plots / "eeg_torus_order_ablation_f1_barplot.png", "fit"),
                t(8.35, 1.55, 3.90, 0.38, "WHAT WAS REMOVED", 12, RED, True),
                *bullet_block(8.35, 2.04, 3.95, ["minor radius r", "fit MSE", "mean surface error", "fraction inside"], 17),
                t(8.35, 4.65, 3.95, 0.35, "INTERPRETATION", 12, RED, True),
                t(8.35, 5.10, 3.98, 1.05, "A quick 15D-to-11D feature ablation, not a refitted lower-order model.", 20, INK, True),
            ],
            f"The complete 15-feature vector reaches {fmt_f1(eeg['ablation_full'])}; removing radius and tube-quality terms gives {fmt_f1(eeg['ablation_no_r'])}. The paired p-value is {eeg['ablation_p']:.4f}, so the decrease is not statistically reliable. This experiment is deliberately modest: it is a feature ablation, not direct measurement of wrong AR-order recovery.",
            "Radius-related geometry helps numerically, but this shortcut does not establish a significant order effect.",
            "The second application asks whether the framework extends beyond sleep to repeated reaching.",
            "Do not label these bars correct-order and incorrect-order model accuracy; both are decoders built from cached geometric features.",
        )
    )

    slides.append(
        standard_slide(
            "A reviewer-requested motor task extends the test across species and behavior",
            "Macaque LFP",
            "core",
            [
                img(0.58, 1.30, 8.35, 5.65, motor_plots / "task_sanity.png", "fit"),
                *metric_card(9.25, 1.48, "Dataset", "2 macaques", RED, 3.25),
                *metric_card(9.25, 2.92, "Converted", f"{counts['motor_unique_lfps']} unique LFPs", MID_RED, 3.25),
                *metric_card(9.25, 4.36, "Main analysis", f"{counts['motor_full_lfps']} six-direction LFPs", BLUE, 3.25),
                t(9.35, 5.96, 3.08, 0.62, "Decode reach direction from the movement-aligned LFP epoch.", 17, INK, True, "center"),
            ],
            "The reviewer specifically asked whether the approach could be applied to periodic trajectories associated with repeated reaching. We used a public motor-cortex LFP dataset from two macaques. The main analysis uses 237 LFP recordings containing all six reach directions and decodes direction from the movement-aligned epoch.",
            "The macaque analysis is a genuinely different task, species, and recording setting.",
            "Before fitting geometry, we made lag parameters trace-specific rather than imposing one global choice.",
            "The documented task is six reach directions crossed with short and long delay; the primary decoder target here is six-way direction, not an invented eight-condition label.",
        )
    )

    dims = emb["dim_counts"]
    dim_text = " | ".join(f"{dim}D: {count}" for dim, count in dims.items())
    slides.append(
        standard_slide(
            "Lag parameters are estimated independently for each LFP trace",
            "Macaque LFP",
            "core",
            [
                img(0.72, 1.30, 5.60, 5.60, motor_plots / "macaque_lfp_pertrace_embedding_figure4c_style_pertrace_tau_dim.png", "fit"),
                t(6.80, 1.55, 5.15, 0.36, "TAU: TEMPORAL SEPARATION", 12, RED, True),
                t(6.80, 2.03, 5.25, 0.98, "First local minimum of average mutual information, with autocorrelation fallbacks.", 21, INK, True),
                t(6.80, 3.30, 5.15, 0.36, "DIMENSION: MODE COUNT", 12, RED, True),
                t(6.80, 3.78, 5.25, 0.98, "Robust PSD peaks after 1/f correction; embedding dimension = 2K + 1, capped at 9D.", 21, INK, True),
                t(6.80, 5.22, 2.45, 0.82, f"Median tau\n{emb['median_tau_ms']:.0f} ms", 23, RED, True, "center", "middle", LIGHT_RED, RED),
                t(9.52, 5.22, 2.65, 0.82, f"Dimensions\n3D to 9D", 23, RED, True, "center", "middle", LIGHT_RED, RED),
                t(6.83, 6.30, 5.30, 0.35, dim_text, 13, MUTED, True, "center"),
            ],
            f"For each unique LFP trace, tau comes from the first local minimum of average mutual information, with autocorrelation fallbacks when necessary. The embedding dimension is derived from robust PSD peaks using two coordinates per oscillatory mode plus one, capped at nine dimensions. The median tau is {emb['median_tau_ms']:.0f} milliseconds, and selected dimensions range from three to nine.",
            "The torus fit is not based on a single arbitrary lag or one global dimension.",
            "With those unsupervised parameters fixed, we decode the six reach directions.",
            "The PSD is used to choose the number of embedding coordinates, not to choose the decoder label or optimize F1.",
        )
    )

    slides.append(
        standard_slide(
            "Torus features decode six reach directions above compact baselines",
            "Macaque LFP",
            "core",
            [
                img(0.62, 1.28, 8.15, 5.70, motor_plots / "direction_movement_feature_f1_barplot_nonlinear_torus_full_6_direction_lfps_pertrace_tau_dim.png", "fit"),
                *metric_card(9.05, 1.55, "Torus features", fmt_f1(mot["torus"]), RED, 3.45),
                *metric_card(9.05, 3.02, "Relevant band: beta", fmt_f1(mot["relevant"]), MID_RED, 3.45),
                *metric_card(9.05, 4.49, "Average PSD", fmt_f1(mot["average_psd"]), "888888", 3.45),
                t(9.20, 6.12, 3.15, 0.38, "Chance F1 approx. 0.17", 14, MUTED, True, "center"),
            ],
            f"For the 237 full six-direction LFPs, torus features reach {fmt_f1(mot['torus'])}. Beta power reaches {fmt_f1(mot['relevant'])}, and Average PSD reaches {fmt_f1(mot['average_psd'])}. The effect size is modest in absolute terms, as expected for single-channel six-way decoding, but it is consistent across a large number of LFP recordings.",
            "Geometry carries reach-direction information beyond beta power and mean PSD.",
            "Matched tests quantify how consistently that advantage appears across LFPs.",
            "Chance is approximately one-sixth, but macro F1 can sit slightly below nominal chance for noisy multiclass predictions; the paired comparisons are the more informative test.",
        )
    )

    stat_body: list[Element] = []
    stat_body.extend(metric_card(0.85, 1.65, "Torus minus beta", "+0.024 F1", RED, 3.55))
    stat_body.extend(metric_card(4.90, 1.65, "Torus minus Average PSD", "+0.022 F1", RED, 3.55))
    stat_body.extend(metric_card(8.95, 1.65, "Matched units", "237 LFPs", BLUE, 3.55))
    stat_body.extend(
        [
            t(1.10, 3.55, 4.70, 0.45, fmt_p(mot["torus_vs_relevant_p_holm"]), 24, RED, True, "center"),
            t(7.55, 3.55, 4.70, 0.45, fmt_p(mot["torus_vs_psd_p_holm"]), 24, RED, True, "center"),
            line(2.00, 4.22, 9.35, 0, GRID, 1.4),
            t(1.10, 4.65, 11.10, 0.74, "The small mean difference is highly consistent across independent LFP recordings.", 26, INK, True, "center"),
            t(1.32, 5.70, 10.70, 0.58, "Paired Wilcoxon signed-rank tests; Holm correction across the three planned comparisons.", 17, MUTED, False, "center"),
        ]
    )
    slides.append(
        standard_slide(
            "Small macaque geometry gains are highly consistent",
            "Macaque LFP",
            "core",
            stat_body,
            f"The average improvements are about 0.024 F1 over beta and 0.022 over Average PSD. Because these are paired within the same 237 LFPs, we can ask whether the sign and magnitude of the difference are consistent. Both survive Holm correction: {fmt_p(mot['torus_vs_relevant_p_holm'])} against beta and {fmt_p(mot['torus_vs_psd_p_holm'])} against Average PSD.",
            "The effect is modest but highly repeatable across LFP traces.",
            "The same qualitative pattern also appears separately in each macaque.",
            "Statistical significance here reflects a large matched sample; pair the p-values with the effect sizes and SDs rather than presenting p-values alone.",
        )
    )

    slides.append(
        standard_slide(
            "The qualitative result replicates separately in both macaques",
            "Macaque LFP",
            "core",
            [
                img(0.70, 1.40, 8.55, 4.90, motor_plots / "by_monkey" / "direction_movement_feature_f1_heatmap_by_monkey_pertrace_tau_dim.png", "fit"),
                t(9.60, 1.78, 2.80, 0.35, "MONKEY M", 12, RED, True),
                t(9.60, 2.25, 2.80, 0.62, "Torus F1\n0.175 +/- 0.034", 22, INK, True),
                t(9.60, 3.28, 2.80, 0.35, "MONKEY T", 12, RED, True),
                t(9.60, 3.75, 2.80, 0.62, "Torus F1\n0.188 +/- 0.028", 22, INK, True),
                t(9.42, 5.18, 3.15, 0.78, "Replication by animal - not train-on-one, test-on-the-other transfer.", 18, RED, True, "center"),
            ],
            "When we summarize Monkey M and Monkey T separately, torus features exceed beta and Average PSD in both animals. Monkey T has higher overall decodability, especially for all-band power, but the compact-baseline geometry advantage is not restricted to one monkey.",
            "Both animals independently support the compact-baseline result.",
            "Finally, combining representations asks whether geometry and power are redundant or complementary.",
            "Call this replication by animal. The decoder is fit and evaluated within each LFP; this is not a leave-one-animal-out transfer analysis.",
        )
    )

    slides.append(
        standard_slide(
            "Geometry complements beta and Average PSD, not all-band power",
            "Macaque LFP",
            "core",
            [
                img(0.62, 1.30, 8.50, 5.65, motor_plots / "additive_torus" / "pooled_additive_torus_comparison_f1_barplot_pertrace_tau_dim.png", "fit"),
                t(9.50, 1.55, 2.82, 0.33, "ADDITIVE EFFECT", 12, RED, True),
                t(9.50, 2.05, 2.82, 0.72, "+0.029 F1\nTorus + beta", 22, INK, True),
                t(9.50, 3.05, 2.82, 0.72, "+0.030 F1\nTorus + Average PSD", 22, INK, True),
                t(9.50, 4.18, 2.82, 0.72, "-0.002 F1\nTorus + all bands", 22, INK, True),
                t(9.42, 5.48, 3.02, 0.82, "Complementarity is strongest when the spectral baseline is compact.", 18, RED, True, "center"),
            ],
            "Concatenating torus features with beta increases F1 by about 0.029, and concatenating them with Average PSD increases F1 by about 0.030. Both effects are strongly significant. But adding torus features to the full five-band vector does not improve the pooled result. The all-band vector already captures much of the task information available to this simple linear decoder.",
            "Geometry contributes nonredundant information to compact spectral summaries, not an unlimited independent signal.",
            "That leads to the balanced conclusion of the rebuttal experiments.",
            "Concatenation increases dimensionality. The fair evidence is the paired out-of-fold F1, not the training score; no improvement over all-band power argues against a generic dimension-only benefit.",
        )
    )

    conclusion_body: list[Element] = []
    claims = [
        ("01", "BEHAVIORAL UTILITY", "Torus features decode sleep stage and reach direction across two species and recording modalities."),
        ("02", "COMPLEMENTARITY", "Geometry consistently beats a relevant single band and Average PSD, and improves them when concatenated."),
        ("03", "CALIBRATED CLAIM", "All-band power remains stronger than torus alone; the contribution is a principled geometric representation, not universal decoder supremacy."),
    ]
    for i, (num, label, text_value) in enumerate(claims):
        y = 1.42 + i * 1.58
        conclusion_body.extend(
            [
                t(0.75, y, 0.68, 0.50, num, 23, RED, True, "center"),
                line(1.58, y + 0.23, 0.65, 0, RED, 2),
                t(2.48, y - 0.03, 3.10, 0.30, label, 12, MUTED, True),
                t(2.48, y + 0.36, 9.55, 0.72, text_value, 20, INK, True),
            ]
        )
    conclusion_body.extend(
        [
            rect(0.72, 6.20, 11.90, 0.60, RED),
            t(0.95, 6.29, 11.45, 0.38, "Next: Rudra - reproducible package, GitHub tutorial, and future applications", 17, WHITE, True, "center"),
        ]
    )
    slides.append(
        standard_slide(
            "The rebuttal turns the geometry into a testable application story",
            "Conclusions",
            "core",
            conclusion_body,
            "The application takeaway has three parts. First, the recovered geometry is behaviorally informative in both sleep and reaching. Second, it contributes information beyond compact spectral summaries. Third, the strongest multiband spectral vector remains competitive or stronger, so our claim is complementarity and interpretability rather than universal decoding superiority. The rebuttal made the paper more precise and more useful.",
            "The geometry is useful, reproducible, and scientifically interpretable - with clearly stated limits.",
            "Rudra will now show how the full workflow is packaged and how a new user can apply it.",
            "If challenged on why publish a representation that does not always win decoding: prediction accuracy and geometric fidelity are distinct objectives; the geometry links oscillatory structure to an interpretable state-space object.",
            "Slow down here. This is the slide the audience should remember.",
        )
    )

    # Backup 1: exact features
    slides.append(
        standard_slide(
            "The torus decoder uses 15 geometric features",
            "Feature definitions",
            "backup",
            [
                t(0.80, 1.38, 3.65, 0.35, "SIZE", 12, RED, True),
                *bullet_block(0.80, 1.86, 3.70, ["Major radius R1", "Major radius R2", "Tube radius r"], 19),
                t(4.80, 1.38, 3.65, 0.35, "FIT QUALITY", 12, RED, True),
                *bullet_block(4.80, 1.86, 3.70, ["Mean-squared error", "Mean surface error", "Fraction inside"], 19),
                t(8.80, 1.38, 3.65, 0.35, "POSE: 9 COMPONENTS", 12, RED, True),
                *bullet_block(8.80, 1.86, 3.70, ["Direction x, y, z", "u-axis x, y, z", "v-axis x, y, z"], 19),
                rect(0.85, 4.50, 11.65, 1.25, LIGHT_RED, RED),
                t(1.18, 4.79, 10.98, 0.65, "3 size + 3 fit-quality + 9 orientation = 15 features per analysis window or trial epoch", 22, INK, True, "center"),
                t(1.10, 6.12, 11.15, 0.45, "The decoder sees fitted parameters, not the raw point cloud.", 17, MUTED, True, "center"),
            ],
            "Use this slide when someone asks exactly what the 15 torus features are. The vector combines scale, fit quality, and pose. It is intentionally compact compared with feeding the entire delay cloud to a flexible model.",
            "The torus representation is a compact, interpretable parameter vector.",
            "Return to the relevant result slide.",
            "Orientation terms are retained in the no-radius ablation because they describe the large fitted ellipse's pose, not local tube thickness.",
        )
    )

    slides.append(
        standard_slide(
            "The reviewer-facing torus fit uses nonlinear surface optimization",
            "Torus fitting",
            "backup",
            [
                rect(0.75, 1.52, 2.55, 1.25, WHITE, GRID),
                t(0.95, 1.78, 2.15, 0.65, "1. Delay-embed\nthe LFP epoch", 19, INK, True, "center"),
                line(3.30, 2.12, 0.55, 0, RED, 2),
                rect(3.85, 1.52, 2.55, 1.25, WHITE, GRID),
                t(4.05, 1.78, 2.15, 0.65, "2. PCA initialize\ncenter and axes", 19, INK, True, "center"),
                line(6.40, 2.12, 0.55, 0, RED, 2),
                rect(6.95, 1.52, 2.55, 1.25, WHITE, GRID),
                t(7.15, 1.78, 2.15, 0.65, "3. Nonlinear\nelliptical-torus fit", 19, INK, True, "center"),
                line(9.50, 2.12, 0.55, 0, RED, 2),
                rect(10.05, 1.52, 2.55, 1.25, WHITE, GRID),
                t(10.25, 1.78, 2.15, 0.65, "4. Export\n15 parameters", 19, INK, True, "center"),
                t(0.90, 3.55, 3.20, 0.35, "OBJECTIVE", 12, RED, True),
                t(0.90, 4.02, 3.25, 1.10, "Minimize mean squared distance from embedded points to the fitted torus surface.", 20, INK, True),
                t(4.65, 3.55, 3.20, 0.35, "CONSTRAINTS", 12, RED, True),
                t(4.65, 4.02, 3.25, 1.10, "Positive tube radius and bounded major radius prevent degenerate fits.", 20, INK, True),
                t(8.40, 3.55, 3.20, 0.35, "COMPUTE", 12, RED, True),
                t(8.40, 4.02, 3.25, 1.10, "Up to 1,200 objective evaluations per trial; cached after fitting.", 20, INK, True),
                t(1.05, 5.80, 11.05, 0.60, "The final macaque results use this slower nonlinear fit, not the earlier PCA/SVD screening proxy.", 20, RED, True, "center"),
            ],
            "The final macaque pipeline initializes from PCA, then refines the full elliptical torus by nonlinear least squares. Max nfev equals 1200, meaning the optimizer can evaluate the residual objective up to 1200 times. The result is cached because this is the expensive stage.",
            "The reported macaque features come from nonlinear fitted geometry.",
            "Return to the feature comparison.",
            "PCA is only an initialization and a projection for embeddings above three dimensions; it is not the final torus parameter estimate.",
        )
    )

    slides.append(
        standard_slide(
            "Spectral baselines range from one scalar to a multiband vector",
            "Spectral definitions",
            "backup",
            [
                t(0.82, 1.40, 3.60, 0.35, "SINGLE-BAND POWER", 12, RED, True),
                t(0.82, 1.92, 3.65, 1.16, "One log-power scalar per epoch. Each band is decoded separately.", 21, INK, True),
                t(4.85, 1.40, 3.60, 0.35, "AVERAGE PSD", 12, RED, True),
                t(4.85, 1.92, 3.65, 1.16, "Welch PSD averaged over the analysis range, then log10 transformed: one scalar.", 21, INK, True),
                t(8.88, 1.40, 3.60, 0.35, "ALL-BAND POWER", 12, RED, True),
                t(8.88, 1.92, 3.65, 1.16, "All band-power values entered jointly as a multifeature vector; they are not averaged.", 21, INK, True),
                rect(0.85, 3.68, 11.60, 1.35, PALE, GRID),
                t(1.15, 3.98, 11.00, 0.72, "EEG bands: delta, theta, alpha, sigma, beta, low gamma | Macaque bands: delta, theta, alpha, beta, low gamma", 18, INK, True, "center"),
                t(1.10, 5.56, 11.15, 0.78, "Relevant band is fixed as delta for sleep and beta for motor LFP. It is not chosen by maximizing F1.", 20, RED, True, "center"),
            ],
            "This slide clarifies the baseline hierarchy. Average PSD is one scalar; a relevant band is one scalar; all-band power is a vector containing every band power jointly. The richer all-band representation is therefore a substantially stronger baseline.",
            "Baseline complexity increases from one scalar to a multiband vector.",
            "Return to the spectral nuance slide.",
            "Average PSD is log10 of mean Welch PSD after detrending and z-scoring; it is not the average of the decoder scores across bands.",
        )
    )

    slides.append(
        standard_slide(
            "Cross-validation uses matched labels and independent units",
            "Decoding protocol",
            "backup",
            [
                rect(0.80, 1.48, 3.55, 3.92, WHITE, GRID),
                t(1.08, 1.78, 2.95, 0.35, "MOUSE EEG", 12, RED, True),
                *bullet_block(1.08, 2.28, 2.95, ["Balance Wake/NREM/REM within hour", "Five stratified folds", "21 hours summarized independently"], 17),
                rect(4.88, 1.48, 3.55, 3.92, WHITE, GRID),
                t(5.16, 1.78, 2.95, 0.35, "MACAQUE LFP", 12, RED, True),
                *bullet_block(5.16, 2.28, 2.95, ["Balance six direction labels within LFP", "Five stratified folds", "237 LFPs summarized independently"], 17),
                rect(8.96, 1.48, 3.55, 3.92, WHITE, GRID),
                t(9.24, 1.78, 2.95, 0.35, "COMPARISON", 12, RED, True),
                *bullet_block(9.24, 2.28, 2.95, ["Same folds for matched features", "Macro F1 from out-of-fold predictions", "Paired Wilcoxon + Holm correction"], 17),
                t(1.05, 5.82, 11.15, 0.60, "No session or LFP is duplicated as an independent observation in the reported summary statistics.", 20, RED, True, "center"),
            ],
            "Each decoding result is cross-validated within its independent reporting unit: hour for EEG and unique LFP for macaque. Features are compared on matched units, then paired statistics are run across those units.",
            "The statistics operate on independent session- or LFP-level F1 values.",
            "Return to the relevant p-value slide.",
            "The existing analysis is not a session-disjoint train/test transfer decoder; describe it as repeated within-unit cross-validation summarized across independent units.",
        )
    )

    slides.append(
        standard_slide(
            "Full EEG heatmap: 21 valid recording hours",
            "EEG detail",
            "backup",
            [img(1.00, 1.28, 11.30, 5.95, eeg_plots / "session_f1_heatmap.png", "fit")],
            "Use this full-size heatmap when the audience wants the individual hourly values. It exposes both the stability and the real session-to-session variability.",
            "The full session matrix is available rather than only a collapsed mean.",
            "Return to the EEG summary.",
            "Hours 1, 2, and 14 are absent because they lacked enough examples from all three classes.",
        )
    )

    slides.append(
        standard_slide(
            "Macaque confusion matrices show structured direction signals",
            "Macaque detail",
            "backup",
            [
                img(0.70, 1.42, 5.85, 5.15, motor_plots / "direction_movement_torus_nonlinear_mean_confusion_pertrace_tau_dim.png", "fit"),
                img(6.78, 1.42, 5.85, 5.15, motor_plots / "direction_movement_beta_mean_confusion_pertrace_tau_dim.png", "fit"),
                t(1.20, 6.55, 4.90, 0.35, "Torus features", 16, RED, True, "center"),
                t(7.28, 6.55, 4.90, 0.35, "Relevant band: beta", 16, MID_RED, True, "center"),
            ],
            "These averaged confusion matrices show that six-way decoding is challenging but not structureless. The torus decoder yields a slightly stronger diagonal than the beta-only decoder across the same LFPs.",
            "The macaque F1 difference corresponds to distributed direction information, not one isolated class.",
            "Return to the main macaque comparison.",
            "A mean confusion matrix averages normalized out-of-fold confusion matrices across LFPs; it is not one pooled classifier.",
        )
    )

    add_rows = mot["additive_significance"]
    add_table_elements: list[Element] = [
        rect(0.75, 1.48, 11.85, 0.62, RED),
        t(0.95, 1.62, 4.60, 0.30, "PAIRED COMPARISON", 13, WHITE, True),
        t(6.00, 1.62, 2.15, 0.30, "DELTA F1", 13, WHITE, True, "center"),
        t(8.45, 1.62, 2.15, 0.30, "HOLM p", 13, WHITE, True, "center"),
        t(10.80, 1.62, 1.35, 0.30, "CALL", 13, WHITE, True, "center"),
    ]
    table_rows = [
        ("Torus + Relevant band vs Relevant band", add_rows["torus_plus_relevant_band vs relevant_band"]),
        ("Torus + Average PSD vs Average PSD", add_rows["torus_plus_average_psd vs average_psd"]),
        ("Torus + All-band power vs All-band power", add_rows["torus_plus_all_band_power vs all_band_power"]),
    ]
    for i, (label, values) in enumerate(table_rows):
        y = 2.12 + i * 1.02
        fill = WHITE if i % 2 == 0 else PALE
        add_table_elements.extend(
            [
                rect(0.75, y, 11.85, 0.92, fill, GRID),
                t(0.98, y + 0.22, 4.72, 0.42, label, 17, INK, True),
                t(6.00, y + 0.22, 2.15, 0.42, f"{values['difference']:+.3f}", 18, INK, True, "center"),
                t(8.30, y + 0.22, 2.45, 0.42, fmt_p(values["p_holm"]).replace("p = ", ""), 16, INK, True, "center"),
                t(10.90, y + 0.22, 1.15, 0.42, "***" if values["p_holm"] < 0.001 else "ns", 18, RED, True, "center"),
            ]
        )
    add_table_elements.append(t(1.00, 5.72, 11.25, 0.72, "The additive result is specific: geometry supplements compact baselines but does not improve the complete multiband vector.", 21, RED, True, "center"))
    slides.append(
        standard_slide(
            "Additive tests isolate complementary geometric information",
            "Additive statistics",
            "backup",
            add_table_elements,
            "This table gives the exact primary additive comparisons. Torus features significantly improve beta and Average PSD. They do not improve all-band power in the pooled analysis.",
            "Complementarity depends on what the spectral baseline already contains.",
            "Return to the additive bar plot.",
            "Holm correction is applied within the pooled additive comparison family shown here.",
        )
    )

    slides.append(
        standard_slide(
            "Per-trace lag selection yields heterogeneous dimensions",
            "Embedding detail",
            "backup",
            [
                img(0.65, 1.35, 7.55, 5.55, motor_plots / "lfp_embedding_parameter_summary_pertrace_tau_dim.png", "fit"),
                *metric_card(8.62, 1.58, "Median tau", f"{emb['median_tau_ms']:.0f} ms", RED, 3.40),
                *metric_card(8.62, 3.05, "Tau IQR", f"{emb['tau_q1_ms']:.0f}-{emb['tau_q3_ms']:.0f} ms", MID_RED, 3.40),
                t(8.75, 4.65, 3.12, 0.35, "DIMENSION COUNTS", 12, RED, True, "center"),
                t(8.75, 5.12, 3.12, 1.02, "\n".join(f"{dim}D: {count} traces" for dim, count in dims.items()), 18, INK, True, "center"),
            ],
            "This diagnostic shows why we abandoned a single dataset-wide tau and dimension. The signal structure varies across electrodes and sessions, so each trace receives unsupervised parameters before decoding.",
            "Heterogeneity in oscillatory content is modeled rather than averaged away.",
            "Return to the per-trace embedding slide.",
            "Although embeddings may be higher than three dimensions, the current 15-feature torus fit projects each cloud to its leading three principal coordinates for a comparable elliptical-torus parameterization.",
        )
    )

    slides.append(
        standard_slide(
            "Persistent homology tests the torus hypothesis rather than forcing it",
            "Topology",
            "backup",
            [
                img(0.65, 1.35, 5.72, 5.42, root / "plots" / "topology" / "session_m0900_0960" / "threshold_sweep.png", "fit"),
                img(6.80, 1.35, 5.72, 5.42, root / "plots" / "topology" / "session_m0900_0960" / "betti_signature_counts.png", "fit"),
                t(0.95, 6.72, 11.48, 0.28, "Observed dominant EEG signature in this probe: (1, 1, 0), not the hypothesized (1, 2, 1).", 15, RED, True, "center"),
            ],
            "This exploratory EEG persistent-homology notebook swept delays, window lengths, and thresholds without hardcoding the answer. In the selected session, the expected two-torus signature was not dominant. That negative result is informative and should not be conflated with the separate fitted-feature decoding analysis.",
            "Topological validation and geometric-feature utility are related but distinct empirical questions.",
            "Return to the limitations or conclusions slide.",
            "Do not use this slide as evidence that every EEG window is a clean topological torus; it demonstrates the diagnostic and its honesty.",
        )
    )

    evidence_rows = [
        ("Specified baselines and error bars", "EEG and macaque relevant-band, Average PSD, and all-band comparisons", "Addressed"),
        ("Stability across recordings", "21 EEG hours; 237 macaque LFPs; separate M/T summaries", "Addressed"),
        ("Repeated motor behavior", "Six-way reach-direction decoding from movement-aligned LFP", "Addressed"),
        ("How much data / lag choice", "Per-trace AMI tau; PSD-informed dimension; parameter distributions", "Quantified"),
        ("Order versus downstream utility", "15D-to-11D no-radius feature ablation", "Partial"),
    ]
    evidence_elements: list[Element] = [
        rect(0.55, 1.35, 12.20, 0.58, RED),
        t(0.75, 1.48, 3.35, 0.28, "REVIEW QUESTION", 12, WHITE, True),
        t(4.38, 1.48, 6.40, 0.28, "NEW EVIDENCE", 12, WHITE, True),
        t(11.05, 1.48, 1.40, 0.28, "STATUS", 12, WHITE, True, "center"),
    ]
    for i, (q, evidence, status) in enumerate(evidence_rows):
        y = 1.94 + i * 0.92
        fill = WHITE if i % 2 == 0 else PALE
        evidence_elements.extend(
            [
                rect(0.55, y, 12.20, 0.84, fill, GRID),
                t(0.75, y + 0.15, 3.30, 0.55, q, 15, INK, True),
                t(4.38, y + 0.15, 6.35, 0.55, evidence, 15, INK),
                t(11.05, y + 0.19, 1.40, 0.42, status, 14, GREEN if status != "Partial" else MID_RED, True, "center"),
            ]
        )
    evidence_elements.append(t(0.90, 6.70, 11.55, 0.28, "The remaining gaps are explicit, which strengthens rather than weakens the revised claim.", 16, RED, True, "center"))
    slides.append(
        standard_slide(
            "The rebuttal converts reviewer concerns into concrete evidence",
            "Reviewer matrix",
            "backup",
            evidence_elements,
            "Use this as a navigation slide during questions. It maps each major concern to a result and marks the order-ablation response as partial rather than overselling it.",
            "Most application concerns now have a direct table, figure, or diagnostic.",
            "Close by returning to the calibrated conclusion.",
            "The remaining theoretical noise/window-placement bounds belong in Kasra's theory discussion rather than this application section.",
        )
    )

    if len(slides) != 25:
        raise ValueError(f"Expected 25 slides, built {len(slides)}")
    return slides


def build_slide_specs(results: dict[str, Any]) -> list[SlideSpec]:
    root = Path(results["root"])
    motor = Path(results["motor_root"])
    eeg_plots = root / "plots" / "summary"
    motor_plots = motor / "plots" / "summary"
    assets = Path(__file__).resolve().parent / "assets"
    counts = results["counts"]
    eeg = results["eeg"]
    mot = results["motor"]

    slides = [
        title_slide(
            "From torus geometry to behavior",
            "Mouse EEG sleep staging and macaque reach-direction decoding",
            "Kasra has established the theoretical link between sustained oscillations and toroidal lag geometry. I will focus on the behavioral tests: first sleep-stage structure and decoding across a full day of mouse EEG, then reach-direction decoding from macaque motor-cortex LFP.",
        ),
        standard_slide(
            "Sleep geometry",
            "Mouse EEG",
            "core",
            [img(1.52, 1.22, 10.30, 5.88, assets / "eeg_sleep_geometry.png", "crop")],
            "We begin with the geometry itself. Each two-second EEG window is delay-embedded and fit with an elliptical torus, summarized here by the two major radii R1 and R2 and the tube radius r. The distributions shift with behavioral state. NREM is especially distinct in R2 and r, while Wake and REM partially overlap. This is the visual reason to ask whether the geometry can decode sleep stage.",
            "The fitted torus changes systematically across Wake, NREM, and REM.",
            "The next question is whether those shifts support class-specific decoding.",
            "These are feature distributions from the original analysis; they motivate decoding but are not themselves a cross-validated performance result.",
            "Point to R2 and r first; they show the clearest separation.",
        ),
        standard_slide(
            "Sleep-stage decoding",
            "Mouse EEG",
            "core",
            [img(0.66, 1.28, 12.02, 5.82, assets / "eeg_class_accuracy.png", "fit")],
            "The left matrix is the original proof-of-concept example. The right matrix is the more conservative rebuttal summary: the mean row-normalized out-of-fold confusion matrix across 21 valid one-hour sessions, with the cell-to-cell standard deviation across hours. NREM is the most reliably identified class. Wake and REM are more often confused with one another, and that structure reproduces across the full-day analysis.",
            "The class structure in the original example persists across independent recording hours.",
            "We next compare geometry against explicit spectral feature choices.",
            "These cells are class-wise recalls, not F1 values. The right panel averages normalized matrices across independently decoded hours; it is not one pooled 24-hour classifier.",
            "Read the diagonal, then mention the Wake-REM confusion pattern.",
        ),
        standard_slide(
            "EEG feature comparison",
            "Mouse EEG",
            "core",
            [img(0.55, 1.23, 12.23, 5.90, eeg_plots / "eeg_sleep_feature_f1_barplot.png", "fit")],
            f"Across {counts['eeg_valid_sessions']} valid hours, the 15 torus features reach {fmt_f1(eeg['torus'])}. They exceed every single-band baseline and the one-scalar Average PSD baseline. Delta is fixed a priori as the physiologically relevant sleep band. The complete all-band vector reaches {fmt_f1(eeg['all_band'])}, slightly above torus alone, so the calibrated claim is that geometry beats compact spectral summaries while the full multiband vector remains a strong baseline.",
            "Geometry outperforms each compact spectral baseline; all-band power is strongest overall.",
            "The hour-by-hour view tests whether those averages are driven by only a few sessions.",
            "Error bars are SD across hours. All-band power is a six-feature vector, not an average of the six band powers.",
            "State the torus result, then immediately give the all-band qualification.",
        ),
        standard_slide(
            "Across-hour stability",
            "Mouse EEG",
            "core",
            [img(3.18, 1.22, 7.00, 5.98, eeg_plots / "session_f1_heatmap.png", "fit")],
            "Each row is one independently decoded recording hour and each column is one representation. Twenty-one of the 24 hours contain enough Wake, NREM, and REM windows for balanced five-fold decoding. The torus column remains consistently high across the day and is usually stronger than any single frequency band, so the mean effect is not produced by one isolated session.",
            "The EEG result repeats across 21 independent one-hour analyses.",
            "We then remove the radius-related information to test which geometric terms matter.",
            "This is within-hour cross-validation repeated across hours, not train-on-one-hour and test-on-another transfer.",
            "Trace the torus column vertically rather than reading individual cells.",
        ),
        standard_slide(
            "Radius-feature ablation",
            "Mouse EEG",
            "core",
            [img(3.15, 1.22, 7.05, 5.98, eeg_plots / "eeg_torus_order_ablation_f1_barplot.png", "fit")],
            f"The complete 15-feature torus vector reaches {fmt_f1(eeg['ablation_full'])}. Removing the tube radius and three tube-quality terms gives {fmt_f1(eeg['ablation_no_r'])}. The paired Wilcoxon p-value is {eeg['ablation_p']:.4f}, so the numerical decrease is not statistically reliable. This is a targeted feature ablation, not a full refit of an intentionally incorrect autoregressive order.",
            "Radius-related terms help numerically, but the ablation difference is not significant.",
            "We now move from spontaneous sleep dynamics to an externally structured motor task.",
            "Do not call this correct-order versus incorrect-order model accuracy; both bars use cached geometric fits and differ only in retained features.",
            "Say 'feature ablation' before interpreting the p-value.",
        ),
        standard_slide(
            "Reaching task",
            "Macaque LFP",
            "core",
            [
                img(0.60, 1.25, 12.13, 4.72, motor_plots / "task_sanity.png", "fit"),
                t(
                    0.78,
                    6.27,
                    11.77,
                    0.46,
                    f"2 macaques   |   6 reach directions   |   movement-aligned LFP   |   {counts['motor_full_lfps']} complete recordings",
                    17,
                    MUTED,
                    True,
                    "center",
                ),
            ],
            "The reviewer asked whether the framework extends to repeated reaching movements. This public motor-cortex dataset contains two macaques performing six reach directions under short and long delay conditions. The sanity panels show balanced task coverage, GO-aligned LFP spectral structure, and movement onset after the GO cue. The primary analysis decodes six-way reach direction from the movement-aligned LFP epoch across 237 recordings containing all directions.",
            "This is a new species, recording modality, and behaviorally instructed task.",
            "Using the same fixed decoder, we compare spectral and torus representations.",
            "The decoder target is reach direction, not task epoch and not a fabricated eight-condition label; both delay types contribute trials.",
            "Use the three panels only to orient the audience, then move on.",
        ),
        standard_slide(
            "Reach-direction decoding",
            "Macaque LFP",
            "core",
            [
                img(
                    0.55,
                    1.23,
                    12.23,
                    5.90,
                    motor_plots
                    / "direction_movement_feature_f1_barplot_nonlinear_torus_full_6_direction_lfps_pertrace_tau_dim.png",
                    "fit",
                )
            ],
            f"Across the {counts['motor_full_lfps']} complete six-direction LFP recordings, torus features reach {fmt_f1(mot['torus'])}. Beta power, fixed a priori as the relevant motor band, reaches {fmt_f1(mot['relevant'])}, and Average PSD reaches {fmt_f1(mot['average_psd'])}. The all-band vector is strongest overall, while torus geometry improves on the compact one-scalar baselines.",
            "Torus geometry carries reach-direction information beyond beta power and Average PSD.",
            "The next slide asks whether the same pattern is visible in each animal separately.",
            "Chance is approximately one-sixth. Report the modest effect size together with the paired consistency across LFPs.",
            "Do not oversell absolute performance; emphasize matched representation comparisons.",
        ),
        standard_slide(
            "Replicates by animal",
            "Macaque LFP",
            "core",
            [
                img(
                    0.55,
                    1.38,
                    12.23,
                    5.48,
                    motor_plots
                    / "by_monkey/direction_movement_feature_f1_heatmap_by_monkey_pertrace_tau_dim.png",
                    "fit",
                )
            ],
            "Separating the data by animal preserves the qualitative result. Torus features exceed beta and Average PSD in both Monkey M and Monkey T. Monkey T has higher overall decodability, particularly for all-band power, but the compact-baseline geometry advantage is not restricted to one animal.",
            "Both animals independently show the compact-baseline geometry advantage.",
            "Finally, concatenation tests whether geometry and spectral power are redundant or complementary.",
            "This is replication by animal, not train-on-one-macaque and test-on-the-other transfer.",
            "Compare rows rather than individual decimal differences.",
        ),
        standard_slide(
            "Complementary geometry",
            "Macaque LFP",
            "core",
            [
                img(
                    0.55,
                    1.23,
                    12.23,
                    5.90,
                    motor_plots
                    / "additive_torus/pooled_additive_torus_comparison_f1_barplot_pertrace_tau_dim.png",
                    "fit",
                )
            ],
            "Concatenating torus geometry with beta power or Average PSD significantly improves those compact baselines. Adding geometry to the already richer all-band vector does not produce a significant gain. This is the most precise conclusion: geometry contributes information that compact spectral summaries miss, while a complete multiband representation already captures much of the decodable signal.",
            "Geometry complements compact spectral baselines, but not the complete all-band vector.",
            "This closes the application story: interpretable geometry, repeatable behavioral information, and calibrated limits.",
            "The significance brackets are paired Wilcoxon tests with Holm correction across matched LFP recordings; error bars are SD across LFPs.",
            "End on complementarity, not on a universal claim of superiority.",
        ),
    ]
    if len(slides) != 10:
        raise ValueError(f"Expected 10 slides, built {len(slides)}")
    return slides


def _rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _text_align(value: str) -> PP_ALIGN:
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[value]


def _anchor(value: str) -> MSO_ANCHOR:
    return {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[value]


def _aspect_fit(source_w: int, source_h: int, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    scale = min(w / source_w, h / source_h)
    out_w = source_w * scale
    out_h = source_h * scale
    return x + (w - out_w) / 2, y + (h - out_h) / 2, out_w, out_h


class PresentationBuilder:
    def __init__(self, slides: Sequence[SlideSpec]):
        self.slides = list(slides)

    def render_pptx(self, output: Path) -> None:
        output.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        blank = prs.slide_layouts[6]

        for index, spec in enumerate(self.slides, start=1):
            slide = prs.slides.add_slide(blank)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = _rgb(WHITE)

            for element in spec.elements:
                self._add_pptx_element(slide, element)
            self._add_pptx_footer(slide, spec, index)

            notes = slide.notes_slide.notes_text_frame
            notes.text = self._notes_text(spec, index)

        prs.save(output)

    def _add_pptx_footer(self, slide: Any, spec: SlideSpec, index: int) -> None:
        footer = slide.shapes.add_textbox(Inches(0.58), Inches(7.10), Inches(11.9), Inches(0.22))
        tf = footer.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"NeuralFieldManifold  |  {spec.section}"
        p.font.name = FONT_NAME
        p.font.size = Pt(8.5)
        p.font.color.rgb = _rgb(MUTED)
        num = slide.shapes.add_textbox(Inches(12.35), Inches(7.08), Inches(0.40), Inches(0.22))
        ntf = num.text_frame
        ntf.clear()
        np_ = ntf.paragraphs[0]
        np_.text = str(index)
        np_.alignment = PP_ALIGN.RIGHT
        np_.font.name = FONT_NAME
        np_.font.size = Pt(8.5)
        np_.font.color.rgb = _rgb(MUTED)

    def _add_pptx_element(self, slide: Any, element: Element) -> None:
        if isinstance(element, RectElement):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(element.x),
                Inches(element.y),
                Inches(element.w),
                Inches(element.h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(element.fill)
            if element.line:
                shape.line.color.rgb = _rgb(element.line)
                shape.line.width = Pt(1)
            else:
                shape.line.fill.background()
        elif isinstance(element, LineElement):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(element.x),
                Inches(element.y),
                Inches(max(element.w, 0.01)),
                Inches(max(element.h, 0.01)),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(element.color)
            shape.line.fill.background()
        elif isinstance(element, TextElement):
            shape = slide.shapes.add_textbox(
                Inches(element.x), Inches(element.y), Inches(element.w), Inches(element.h)
            )
            if element.fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(element.fill)
            else:
                shape.fill.background()
            if element.line:
                shape.line.color.rgb = _rgb(element.line)
            else:
                shape.line.fill.background()
            tf = shape.text_frame
            tf.clear()
            tf.margin_left = Inches(element.margin)
            tf.margin_right = Inches(element.margin)
            tf.margin_top = Inches(element.margin * 0.75)
            tf.margin_bottom = Inches(element.margin * 0.75)
            tf.vertical_anchor = _anchor(element.valign)
            tf.word_wrap = True
            lines = element.text.split("\n")
            for idx, text_line in enumerate(lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = text_line
                p.alignment = _text_align(element.align)
                p.font.name = FONT_NAME
                p.font.size = Pt(element.size)
                p.font.bold = element.bold
                p.font.color.rgb = _rgb(element.color)
                p.space_after = Pt(0)
                p.space_before = Pt(0)
        elif isinstance(element, ImageElement):
            if not element.path.exists():
                raise FileNotFoundError(element.path)
            with Image.open(element.path) as source:
                if element.mode == "fit":
                    x, y, w, h = _aspect_fit(source.width, source.height, element.x, element.y, element.w, element.h)
                    slide.shapes.add_picture(str(element.path), Inches(x), Inches(y), Inches(w), Inches(h))
                else:
                    target_ratio = element.w / element.h
                    source_ratio = source.width / source.height
                    if source_ratio > target_ratio:
                        crop = (1 - target_ratio / source_ratio) / 2
                        pic = slide.shapes.add_picture(str(element.path), Inches(element.x), Inches(element.y), Inches(element.w), Inches(element.h))
                        pic.crop_left = crop
                        pic.crop_right = crop
                    else:
                        crop = (1 - source_ratio / target_ratio) / 2
                        pic = slide.shapes.add_picture(str(element.path), Inches(element.x), Inches(element.y), Inches(element.w), Inches(element.h))
                        pic.crop_top = crop
                        pic.crop_bottom = crop
        else:
            raise TypeError(element)

    def render_previews(self, output_dir: Path) -> list[Path]:
        output_dir.mkdir(parents=True, exist_ok=True)
        paths: list[Path] = []
        for index, spec in enumerate(self.slides, start=1):
            canvas = Image.new("RGB", (PREVIEW_W, PREVIEW_H), "white")
            draw = ImageDraw.Draw(canvas)
            for element in spec.elements:
                self._draw_preview_element(canvas, draw, element)
            self._draw_preview_footer(draw, spec, index)
            path = output_dir / f"slide_{index:02d}.png"
            canvas.save(path, quality=95)
            paths.append(path)
        return paths

    def _coords(self, x: float, y: float, w: float, h: float) -> tuple[int, int, int, int]:
        return (
            round(x / SLIDE_W * PREVIEW_W),
            round(y / SLIDE_H * PREVIEW_H),
            round((x + w) / SLIDE_W * PREVIEW_W),
            round((y + h) / SLIDE_H * PREVIEW_H),
        )

    def _font(self, size: float, bold: bool = False) -> ImageFont.FreeTypeFont:
        path = PREVIEW_FONT_BOLD if bold else PREVIEW_FONT
        return ImageFont.truetype(path, max(8, round(size * 1.33)))

    def _draw_preview_element(self, canvas: Image.Image, draw: ImageDraw.ImageDraw, element: Element) -> None:
        box = self._coords(element.x, element.y, element.w, element.h)
        if isinstance(element, RectElement):
            draw.rectangle(box, fill=f"#{element.fill}", outline=f"#{element.line}" if element.line else None, width=2)
        elif isinstance(element, LineElement):
            x1, y1, x2, y2 = box
            draw.line((x1, y1, x2, y2), fill=f"#{element.color}", width=max(1, round(element.width * 1.3)))
        elif isinstance(element, TextElement):
            if element.fill:
                draw.rectangle(box, fill=f"#{element.fill}", outline=f"#{element.line}" if element.line else None, width=2)
            font = self._font(element.size, element.bold)
            self._draw_wrapped_text(draw, box, element.text, font, f"#{element.color}", element.align, element.valign, element.margin)
        elif isinstance(element, ImageElement):
            with Image.open(element.path).convert("RGB") as source:
                x1, y1, x2, y2 = box
                target_w, target_h = x2 - x1, y2 - y1
                if element.mode == "fit":
                    source.thumbnail((target_w, target_h), Image.Resampling.LANCZOS)
                    paste_x = x1 + (target_w - source.width) // 2
                    paste_y = y1 + (target_h - source.height) // 2
                    canvas.paste(source, (paste_x, paste_y))
                else:
                    source_ratio = source.width / source.height
                    target_ratio = target_w / target_h
                    if source_ratio > target_ratio:
                        crop_w = round(source.height * target_ratio)
                        left = (source.width - crop_w) // 2
                        source = source.crop((left, 0, left + crop_w, source.height))
                    else:
                        crop_h = round(source.width / target_ratio)
                        top = (source.height - crop_h) // 2
                        source = source.crop((0, top, source.width, top + crop_h))
                    source = source.resize((target_w, target_h), Image.Resampling.LANCZOS)
                    canvas.paste(source, (x1, y1))
        else:
            raise TypeError(element)

    def _draw_wrapped_text(
        self,
        draw: ImageDraw.ImageDraw,
        box: tuple[int, int, int, int],
        text_value: str,
        font: ImageFont.FreeTypeFont,
        fill: str,
        align: str,
        valign: str,
        margin: float,
    ) -> None:
        x1, y1, x2, y2 = box
        margin_px = round(margin / SLIDE_W * PREVIEW_W)
        width = max(5, x2 - x1 - 2 * margin_px)
        paragraphs: list[str] = []
        for raw_line in text_value.split("\n"):
            words = raw_line.split()
            if not words:
                paragraphs.append("")
                continue
            current = words[0]
            for word in words[1:]:
                trial = f"{current} {word}"
                if draw.textbbox((0, 0), trial, font=font)[2] <= width:
                    current = trial
                else:
                    paragraphs.append(current)
                    current = word
            paragraphs.append(current)
        line_gap = max(2, round(font.size * 0.10))
        heights = [draw.textbbox((0, 0), value or " ", font=font)[3] for value in paragraphs]
        total_h = sum(heights) + line_gap * max(0, len(paragraphs) - 1)
        if valign == "middle":
            y = y1 + (y2 - y1 - total_h) / 2
        elif valign == "bottom":
            y = y2 - total_h - margin_px
        else:
            y = y1 + margin_px
        for value, height in zip(paragraphs, heights):
            bbox = draw.textbbox((0, 0), value or " ", font=font)
            text_w = bbox[2] - bbox[0]
            if align == "center":
                x = x1 + (x2 - x1 - text_w) / 2
            elif align == "right":
                x = x2 - margin_px - text_w
            else:
                x = x1 + margin_px
            draw.text((x, y), value, font=font, fill=fill)
            y += height + line_gap

    def _draw_preview_footer(self, draw: ImageDraw.ImageDraw, spec: SlideSpec, index: int) -> None:
        font = self._font(8.5)
        draw.text((68, 855), f"NeuralFieldManifold  |  {spec.section}", font=font, fill=f"#{MUTED}")
        text_value = str(index)
        text_w = draw.textbbox((0, 0), text_value, font=font)[2]
        draw.text((1530 - text_w, 855), text_value, font=font, fill=f"#{MUTED}")

    @staticmethod
    def _notes_text(spec: SlideSpec, index: int) -> str:
        return (
            f"SLIDE {index}: {spec.title}\n\n"
            f"SCRIPT\n{spec.script}\n\n"
            f"TAKEAWAY\n{spec.takeaway}\n\n"
            f"TRANSITION\n{spec.transition}\n\n"
            f"IF CHALLENGED\n{spec.challenge_response}\n\n"
            f"DELIVERY CUE\n{spec.speaker_cue or 'Keep the slide conversational and do not read every label.'}"
        )


def write_notes(slides: Sequence[SlideSpec], output: Path) -> None:
    lines = ["# NeuralFieldManifold Rebuttal Applications - Speaker Notes", ""]
    for index, slide in enumerate(slides, start=1):
        prefix = "BACKUP - " if slide.kind == "backup" else ""
        lines.extend(
            [
                f"## Slide {index}: {prefix}{slide.title}",
                "",
                "**Script**",
                "",
                slide.script,
                "",
                f"**Takeaway:** {slide.takeaway}",
                "",
                f"**Transition:** {slide.transition}",
                "",
                f"**If challenged:** {slide.challenge_response}",
                "",
                f"**Delivery cue:** {slide.speaker_cue or 'Keep the slide conversational and do not read every label.'}",
                "",
            ]
        )
    output.write_text("\n".join(lines), encoding="utf-8")


def write_cheat_sheet(results: dict[str, Any], output: Path) -> None:
    eeg = results["eeg"]
    mot = results["motor"]
    text_value = f"""# NeuralFieldManifold Application Section - Rehearsal Sheet

## Opening

"Kasra established why sustained oscillatory modes generate toroidal lag geometry. I will ask whether that recovered geometry contains behaviorally useful information."

## Six Numbers To Know

| Result | F1 |
|---|---:|
| EEG torus | {fmt_f1(eeg['torus'])} |
| EEG delta | {fmt_f1(eeg['relevant'])} |
| EEG Average PSD | {fmt_f1(eeg['average_psd'])} |
| Macaque torus | {fmt_f1(mot['torus'])} |
| Macaque beta | {fmt_f1(mot['relevant'])} |
| Macaque Average PSD | {fmt_f1(mot['average_psd'])} |

EEG paired tests: torus vs delta and torus vs Average PSD, both {fmt_p(eeg['torus_vs_relevant_p_holm'])}.  
Macaque paired tests: torus vs beta {fmt_p(mot['torus_vs_relevant_p_holm'])}; torus vs Average PSD {fmt_p(mot['torus_vs_psd_p_holm'])}.

## Three Qualifications

1. All-band power is a multifeature vector and is stronger than torus alone in both principal datasets.
2. The Monkey M/T analysis is replication by animal, not train-on-one-animal/test-on-the-other transfer.
3. The EEG 15D-versus-11D comparison is a no-radius feature ablation, not a full recovered-order experiment.

## Likely Questions

**Why F1?** Macro F1 weights all sleep stages or reach directions equally; labels are balanced before five-fold LDA.

**What are the error bars?** Standard deviation across 21 EEG sessions or 237 unique full-six-direction LFP recordings.

**What is Average PSD?** One scalar: log10 of Welch PSD averaged across the analysis frequency range after detrending and z-scoring.

**How were macaque tau and dimension chosen?** Tau from average mutual information; dimension from robust 1/f-corrected PSD peaks using 2K+1, capped at 9D. Median tau = {results['embedding']['median_tau_ms']:.0f} ms.

**What is the strongest defensible claim?** Torus geometry carries task information beyond compact spectral summaries and adds to beta or Average PSD, but does not universally outperform a complete multiband representation.

## Closing

"The rebuttal turns the theory into a testable application story: useful geometry, reproducible effects, and calibrated claims. Rudra will now show how the workflow is packaged for new datasets."
"""
    output.write_text(text_value, encoding="utf-8")


def make_contact_sheet(previews: Sequence[Path], output: Path, columns: int = 5) -> None:
    thumb_w, thumb_h = 320, 180
    rows = math.ceil(len(previews) / columns)
    sheet = Image.new("RGB", (columns * thumb_w, rows * thumb_h), "white")
    for index, path in enumerate(previews):
        with Image.open(path).convert("RGB") as source:
            source = source.resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
            x = (index % columns) * thumb_w
            y = (index // columns) * thumb_h
            sheet.paste(source, (x, y))
    sheet.save(output, quality=95)


def previews_to_pdf(previews: Sequence[Path], output: Path) -> None:
    pages = [Image.open(path).convert("RGB") for path in previews]
    try:
        pages[0].save(output, save_all=True, append_images=pages[1:], resolution=144)
    finally:
        for page in pages:
            page.close()


def validate_pptx(path: Path, expected_slides: int) -> dict[str, Any]:
    prs = Presentation(path)
    if len(prs.slides) != expected_slides:
        raise ValueError(f"PPTX slide count {len(prs.slides)} != {expected_slides}")
    titles = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text.strip()]
        if not texts:
            raise ValueError(f"Slide {index} has no text")
        titles.append(texts[0])
    with zipfile.ZipFile(path) as archive:
        broken = [name for name in archive.namelist() if name.endswith(".rels") and archive.getinfo(name).file_size == 0]
        if broken:
            raise ValueError(f"Empty relationship files: {broken}")
    return {"slide_count": len(prs.slides), "titles": titles, "file_size": path.stat().st_size}


def main() -> None:
    presentation_dir = Path(__file__).resolve().parent
    root = presentation_dir.parent
    results = load_results(root)
    validate_claims(results)
    render_eeg_class_accuracy_figure(
        root, presentation_dir / "assets" / "eeg_class_accuracy.png"
    )
    slides = build_slide_specs(results)
    builder = PresentationBuilder(slides)

    pptx_path = presentation_dir / "NeuralFieldManifold_rebuttal_applications.pptx"
    notes_path = presentation_dir / "NeuralFieldManifold_rebuttal_applications_notes.md"
    cheat_path = presentation_dir / "NeuralFieldManifold_rebuttal_applications_cheat_sheet.md"
    pdf_path = presentation_dir / "NeuralFieldManifold_rebuttal_applications_preview.pdf"
    preview_dir = presentation_dir / "previews"

    builder.render_pptx(pptx_path)
    previews = builder.render_previews(preview_dir)
    keep_previews = set(previews)
    for stale in preview_dir.glob("slide_*.png"):
        if stale not in keep_previews:
            stale.unlink()
    make_contact_sheet(previews, preview_dir / "contact_sheet.png")
    previews_to_pdf(previews, pdf_path)
    write_notes(slides, notes_path)
    write_cheat_sheet(results, cheat_path)
    pptx_validation = validate_pptx(pptx_path, len(slides))

    manifest = {
        "title": "NeuralFieldManifold Rebuttal Applications",
        "core_slides": sum(slide.kind == "core" for slide in slides),
        "backup_slides": sum(slide.kind == "backup" for slide in slides),
        "pptx": str(pptx_path),
        "preview_pdf": str(pdf_path),
        "contact_sheet": str(preview_dir / "contact_sheet.png"),
        "speaker_notes": str(notes_path),
        "cheat_sheet": str(cheat_path),
        "pptx_validation": pptx_validation,
        "source_claims": {
            "eeg_torus_f1": eeg_value(results, "torus"),
            "motor_torus_f1": motor_value(results, "torus"),
        },
    }
    (presentation_dir / "presentation_manifest.json").write_text(
        json.dumps(manifest, indent=2), encoding="utf-8"
    )
    print(json.dumps(manifest, indent=2))


def eeg_value(results: dict[str, Any], key: str) -> float:
    return float(results["eeg"][key]["mean_f1"])


def motor_value(results: dict[str, Any], key: str) -> float:
    return float(results["motor"][key]["mean_f1"])


if __name__ == "__main__":
    main()
